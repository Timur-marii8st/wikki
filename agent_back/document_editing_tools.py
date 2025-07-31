import os
import asyncio
from typing import List, Dict, Any

# --- PDF ---
from openpyxl.worksheet.worksheet import Worksheet
from pypdf import PdfWriter, PdfReader
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas

# --- Word (DOCX) ---
from docx import Document

# --- PowerPoint (PPTX) ---
from pptx import Presentation
from pptx.shapes.autoshape import Shape # Импортируем Shape
from pptx.enum.shapes import PP_PLACEHOLDER

# --- Excel (XLSX) ---
from openpyxl import Workbook, load_workbook
from openpyxl.worksheet.worksheet import Worksheet 

# --- ASYNC HELPER ---

async def _run_in_executor(func, *args) -> Any:
    """Запускает блокирующую функцию в отдельном потоке, чтобы не блокировать event loop."""
    return await asyncio.get_running_loop().run_in_executor(None, func, *args)

# ===== ASYNC PDF FUNCTIONS =====

async def create_pdf_from_text(file_path: str, text_content: str, title: str = "Документ", author: str = "AI Agent") -> Dict[str, Any]:
    """
    Асинхронно создает новый PDF-документ из заданного текста.
    Возвращает словарь со статусом операции.
    """
    def worker():
        # Внутренняя функция-воркер со всей логикой
        try:
            c = canvas.Canvas(file_path, pagesize=letter)
            c.setTitle(title)
            c.setAuthor(author)
            
            width, height = letter
            left_margin, right_margin = inch, width - inch
            top_margin, bottom_margin = height - inch, inch
            line_height = 14 
            
            textobject = c.beginText(left_margin, top_margin)
            textobject.setFont("Helvetica", 12)
            
            lines = text_content.split('\n')
            
            for line in lines:
                words = line.split(' ')
                current_line_segment = []
                for word in words:
                    if c.stringWidth(' '.join(current_line_segment + [word]), "Helvetica", 12) < (right_margin - left_margin):
                        current_line_segment.append(word)
                    else:
                        textobject.textLine(' '.join(current_line_segment))
                        current_line_segment = [word]
                textobject.textLine(' '.join(current_line_segment))
            
            c.drawText(textobject)
            c.save()
            return {"status": "success", "message": f"PDF файл '{file_path}' успешно создан."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при создании PDF файла '{file_path}': {e}"}
            
    return await _run_in_executor(worker)


async def merge_pdfs(output_path: str, *input_paths: str) -> Dict[str, Any]:
    """Асинхронно объединяет несколько PDF-файлов в один."""
    def worker():
        try:
            if not input_paths:
                return {"status": "error", "message": "Не указаны входные PDF-файлы для объединения."}
            
            writer = PdfWriter()
            for path in input_paths:
                if not os.path.isfile(path):
                    return {"status": "error", "message": f"Входной файл PDF '{path}' не найден."}
                reader = PdfReader(path)
                for page in reader.pages:
                    writer.add_page(page)
            
            with open(output_path, "wb") as f:
                writer.write(f)
            
            return {"status": "success", "message": f"PDF файлы успешно объединены в '{output_path}'."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при объединении PDF файлов: {e}"}

    return await _run_in_executor(worker)


async def extract_pdf_pages(input_path: str, output_path: str, pages_to_extract: List[int]) -> Dict[str, Any]:
    """Асинхронно извлекает указанные страницы из PDF-файла."""
    def worker():
        try:
            if not os.path.isfile(input_path):
                return {"status": "error", "message": f"Исходный файл PDF '{input_path}' не найден."}
            
            reader = PdfReader(input_path)
            writer = PdfWriter()
            
            for page_num in pages_to_extract:
                if 0 <= page_num < len(reader.pages):
                    writer.add_page(reader.pages[page_num])
                else:
                    return {"status": "error", "message": f"Страница {page_num} выходит за пределы диапазона (0-{len(reader.pages)-1})."}
            
            if not writer.pages:
                return {"status": "warning", "message": "Не удалось извлечь ни одной страницы."}

            with open(output_path, "wb") as f:
                writer.write(f)
            
            return {"status": "success", "message": f"Страницы {pages_to_extract} успешно извлечены в '{output_path}'."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при извлечении страниц из PDF: {e}"}
    
    return await _run_in_executor(worker)


# ===== ASYNC Word (DOCX) FUNCTIONS =====

async def create_docx(file_path: str, content: str = "") -> Dict[str, Any]:
    """Асинхронно создает новый документ Word (.docx)."""
    def worker():
        try:
            document = Document()
            if content:
                document.add_paragraph(content)
            document.save(file_path)
            return {"status": "success", "message": f"DOCX файл '{file_path}' успешно создан."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при создании DOCX файла: {e}"}
    
    return await _run_in_executor(worker)


async def add_text_to_docx(file_path: str, text_to_add: str) -> Dict[str, Any]:
    """Асинхронно добавляет новый абзац в DOCX-документ."""
    def worker():
        try:
            if not os.path.exists(file_path):
                return {"status": "error", "message": f"DOCX файл '{file_path}' не найден."}
            
            document = Document(file_path)
            document.add_paragraph(text_to_add)
            document.save(file_path)
            return {"status": "success", "message": f"Текст успешно добавлен в '{file_path}'."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при добавлении текста в DOCX: {e}"}

    return await _run_in_executor(worker)


async def replace_text_in_docx(file_path: str, old_text: str, new_text: str) -> Dict[str, Any]:
    """Асинхронно заменяет текст в DOCX-документе."""
    def worker():
        try:
            if not os.path.exists(file_path):
                return {"status": "error", "message": f"DOCX файл '{file_path}' не найден."}
            
            document = Document(file_path)
            replaced_count = 0
            for p in document.paragraphs:
                if old_text in p.text:
                    inline = p.runs
                    # Замена текста с сохранением форматирования
                    for i in range(len(inline)):
                        if old_text in inline[i].text:
                            text = inline[i].text.replace(old_text, new_text)
                            inline[i].text = text
                            replaced_count += 1
            
            # ... (аналогичная логика для таблиц)

            document.save(file_path)
            return {"status": "success", "message": f"Успешно заменено {replaced_count} вхождений.", "replaced_count": replaced_count}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при замене текста в DOCX: {e}"}

    return await _run_in_executor(worker)


# ===== ASYNC PowerPoint (PPTX) FUNCTIONS =====

async def create_pptx(file_path: str, title: str = "", subtitle: str = "") -> Dict[str, Any]:
    """Асинхронно создает новую презентацию PowerPoint (.pptx)."""
    def worker():
        try:
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
             # Безопасная проверка и присвоение заголовка
            if title:
                # Проверяем, существует ли shapes.title
                if slide.shapes.title:
                    slide.shapes.title.text = title
                else:
                    # Опционально: добавить текстовое поле, если заголовка нет
                    # Или просто пропустить, если заголовок не является обязательным
                    print("Предупреждение: Заголовок слайда не найден в текущем макете.")

            # Safely setting the subtitle
            if subtitle:
                if len(slide.placeholders) > 1:
                    potential_subtitle_placeholder = slide.placeholders[1]

                    # Проверяем, что это объект, который может иметь текстовый фрейм
                    # PP_PLACEHOLDER и Shape (AutoShape) - это основные типы, которые его имеют
                    # Если вы уверены, что это всегда Placeholder, можно проверить только Placeholder
                    if isinstance(potential_subtitle_placeholder, (Shape, PP_PLACEHOLDER)):
                        if potential_subtitle_placeholder.has_text_frame:
                            potential_subtitle_placeholder.text_frame.text = subtitle
                        else:
                            print(f"Предупреждение: Placeholder по индексу 1 не имеет текстового фрейма для подзаголовка.")
                    else:
                        print(f"Предупреждение: Объект по индексу 1 не является Shape или Placeholder, не может быть подзаголовком.")
                else:
                    print("Предупреждение: Подзаголовок слайда не найден в текущем макете.")

            
            prs.save(file_path)
            return {"status": "success", "message": f"PPTX файл '{file_path}' успешно создан."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при создании PPTX файла: {e}"}

    return await _run_in_executor(worker)


async def add_slide_to_pptx(file_path: str, layout_index: int = 1, title: str = "", content: str = "") -> Dict[str, Any]:
    """Асинхронно добавляет новый слайд в PPTX презентацию."""
    def worker():
        try:
            if not os.path.exists(file_path):
                return {"status": "error", "message": f"PPTX файл '{file_path}' не найден."}
            
            prs = Presentation(file_path)
            if not (0 <= layout_index < len(prs.slide_layouts)):
                return {"status": "error", "message": f"Индекс макета {layout_index} вне диапазона."}

            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            
            if title and slide.shapes.title:
                slide.shapes.title.text = title
            
            # Поиск основного текстового поля и присвоение контента
            for shape in slide.placeholders:
                # 1. Убедимся, что это Shape или Placeholder (для text_frame)
                # 2. Проверяем, что у него есть текстовый фрейм
                # 3. Проверяем, что это не заголовок (используем type вместо is_title)
                if isinstance(shape, (Shape, PP_PLACEHOLDER)) and shape.has_text_frame:
                    if shape.is_placeholder and shape.placeholder_format.type != PP_PLACEHOLDER.TITLE:
                        shape.text_frame.text = content
                        break # Присвоили контент и выходим
            
            prs.save(file_path)
            return {"status": "success", "message": f"Новый слайд успешно добавлен в '{file_path}'."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при добавлении слайда в PPTX: {e}"}

    return await _run_in_executor(worker)


# ===== ASYNC Excel (XLSX) FUNCTIONS =====

async def create_xlsx(file_path: str, sheet_name: str = "Sheet1") -> Dict[str, Any]:
    """Асинхронно создает новый документ Excel (.xlsx)."""
    def worker():
        try:
            wb = Workbook()
            ws = wb.active
            assert isinstance(ws, Worksheet), "Active worksheet is not a Worksheet instance."
            ws.title = sheet_name
            wb.save(file_path)
            return {"status": "success", "message": f"XLSX файл '{file_path}' создан с листом '{sheet_name}'."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при создании XLSX файла: {e}"}

    return await _run_in_executor(worker)


async def write_to_excel_cell(file_path: str, sheet_name: str, cell_address: str, value: Any) -> Dict[str, Any]:
    """Асинхронно записывает значение в ячейку Excel."""
    def worker():
        try:
            if not os.path.exists(file_path): return {"status": "error", "message": f"Файл '{file_path}' не найден."}
            wb = load_workbook(file_path)
            if sheet_name not in wb.sheetnames: return {"status": "error", "message": f"Лист '{sheet_name}' не найден."}
            
            ws = wb[sheet_name]
            ws[cell_address] = value
            wb.save(file_path)
            return {"status": "success", "message": f"Значение '{value}' записано в ячейку {cell_address}."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при записи в ячейку Excel: {e}"}
            
    return await _run_in_executor(worker)


async def read_from_excel_cell(file_path: str, sheet_name: str, cell_address: str) -> Dict[str, Any]:
    """Асинхронно читает значение из ячейки Excel."""
    def worker():
        try:
            if not os.path.exists(file_path): return {"status": "error", "message": f"Файл '{file_path}' не найден."}
            wb = load_workbook(file_path)
            if sheet_name not in wb.sheetnames: return {"status": "error", "message": f"Лист '{sheet_name}' не найден."}
            
            ws = wb[sheet_name]
            value = ws[cell_address].value
            return {"status": "success", "message": f"Значение из ячейки {cell_address} успешно прочитано.", "data": value}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при чтении из ячейки Excel: {e}", "data": None}

    return await _run_in_executor(worker)


async def write_data_to_excel_range(file_path: str, sheet_name: str, start_cell: str, data: List[List[Any]]) -> Dict[str, Any]:
    """Асинхронно записывает 2D-массив данных в Excel."""
    def worker():
        try:
            if not os.path.exists(file_path): return {"status": "error", "message": f"Файл '{file_path}' не найден."}
            wb = load_workbook(file_path)
            if sheet_name not in wb.sheetnames: return {"status": "error", "message": f"Лист '{sheet_name}' не найден."}
            
            ws = wb[sheet_name]
            start_row = ws[start_cell].row
            start_col = ws[start_cell].column
            
            for r_idx, row_data in enumerate(data):
                for c_idx, cell_value in enumerate(row_data):
                    ws.cell(row=start_row + r_idx, column=start_col + c_idx, value=cell_value)
            
            wb.save(file_path)
            return {"status": "success", "message": f"Данные успешно записаны в диапазон начиная с {start_cell}."}
        except Exception as e:
            return {"status": "error", "message": f"Ошибка при записи диапазона в Excel: {e}"}

    return await _run_in_executor(worker)