import os

# Использование LangChain DocumentLoaders - это наиболее рекомендуемый способ
from langchain_community.document_loaders import (
    PyPDFLoader,
    UnstructuredFileLoader,
    TextLoader,
    CSVLoader,
    EvernoteLoader,
    NotebookLoader,
)

import docx 
import pptx
import openpyxl

def read_document(file_path: str) -> str:
    """
    Читает содержимое документа с использованием соответствующих загрузчиков LangChain.
    Эта функция является универсальным интерфейсом.
    """
    if not os.path.exists(file_path):
        return f"Ошибка: Файл '{file_path}' не найден."
    
    # Определяем тип файла по расширению
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    try:
        if file_extension == '.pdf':
            loader = PyPDFLoader(file_path)
        elif file_extension in ('.txt', '.md', '.log', '.py', '.java', '.c', '.cpp', '.h'):
            loader = TextLoader(file_path)
        elif file_extension == '.csv':
            loader = CSVLoader(file_path)
        elif file_extension in ('.docx', '.pptx', '.xlsx', '.odt', '.rtf', '.epub', '.tex', '.html', '.xml'):
            loader = UnstructuredFileLoader(file_path)
        else:
            try:
                # Попробуем прочитать как обычный текст, если формат неизвестен
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except Exception:
                return f"Ошибка: Неизвестный формат файла '{file_extension}' и не удалось прочитать как обычный текст."

        docs = loader.load()
        # Объединяем содержимое всех страниц/частей в одну строку
        full_text = "\n".join([doc.page_content for doc in docs])
        return full_text
    except Exception as e:
        return f"Ошибка при чтении файла '{file_path}' ({file_extension}): {e}"

# --- Специализированные функции для большей точности (если UnstructuredFileLoader не подходит) ---

def read_docx(file_path: str) -> str:
    """
    Читает текст из файла .docx. Требует `pip install python-docx`.
    """
    if not os.path.exists(file_path):
        return f"Ошибка: Файл '{file_path}' не найден."
    try:
        document = docx.Document(file_path)
        full_text = []
        for para in document.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        return f"Ошибка при чтении DOCX файла '{file_path}': {e}"

def read_pptx(file_path: str) -> str:
    """
    Читает текст из файла .pptx (только текстовые поля). Требует `pip install python-pptx`.
    """
    if not os.path.exists(file_path):
        return f"Ошибка: Файл '{file_path}' не найден."
    try:
        presentation = pptx.Presentation(file_path)
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        return f"Ошибка при чтении PPTX файла '{file_path}': {e}"

def read_xlsx(file_path: str) -> str:
    """
    Читает данные из файла .xlsx (все листы). Требует `pip install openpyxl`.
    """
    if not os.path.exists(file_path):
        return f"Ошибка: Файл '{file_path}' не найден."
    try:
        workbook = openpyxl.load_workbook(file_path)
        full_text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            full_text.append(f"--- Лист: {sheet_name} ---")
            for row in sheet.iter_rows():
                row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                full_text.append("\t".join(row_values)) # Используем табуляцию для разделения колонок
        return "\n".join(full_text)
    except Exception as e:
        return f"Ошибка при чтении XLSX файла '{file_path}': {e}"

def read_plain_text(file_path: str) -> str:
    """
    Читает обычный текстовый файл (.txt, .log, .md, .py и т.д.).
    """
    if not os.path.exists(file_path):
        return f"Ошибка: Файл '{file_path}' не найден."
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"Ошибка при чтении обычного текстового файла '{file_path}': {e}"
